import streamlit as st
import docx
from pptx import Presentation
import pdfplumber
import openai
import json

# Веб-парақшаның негізгі баптаулары
st.set_page_config(
    page_title="ҰБТ Материалдарын Тексеру Жүйесі",
    page_icon="🎓",
    layout="wide"
)

# Дизайн стилі
st.markdown("""
    <style>
    .main-header { font-size: 2.2rem; color: #1E3A8A; font-weight: 700; margin-bottom: 0.5rem; }
    .sub-header { font-size: 1.1rem; color: #4B5563; margin-bottom: 2rem; }
    </style>
""", unsafe_allow_html=True)

st.markdown('<div class="main-header">🎓 ҰБТ Материалдарын Автоматты Тексеру Жүйесі</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Word (.docx), PPTX (.pptx) және PDF (.pdf) файлдарындағы қателерді AI методист арқылы табу</div>', unsafe_allow_html=True)

# Сүйер мәзірі (Sidebar) - Баптаулар
with st.sidebar:
    st.header("⚙️ Баптаулар")
    api_key = st.text_input("OpenAI API Key білдіріңіз:", type="password", help="https://platform.openai.com сайтынан алынған API кілт")
    selected_model = st.selectbox("AI Моделі:", ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo"])
    st.markdown("---")
    st.markdown("### 📋 Тексеру бағыттары:")
    check_ortho = st.checkbox("Орфография & Грамматика", value=True)
    check_tech = st.checkbox("Техникалық формат (A-E жауаптар, нөмірлеу)", value=True)
    check_content = st.checkbox("Мазмұндық & Логикалық сәйкестік", value=True)

# Word файлын оқу функциясы
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text_blocks = []
    for i, p in enumerate(doc.paragraphs):
        if p.text.strip():
            text_blocks.append(f"[Абзац {i+1}]: {p.text.strip()}")
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
            if row_text:
                text_blocks.append(f"[Кесте {t_idx+1}, Жол {r_idx+1}]: {row_text}")
    return "\n".join(text_blocks)

# PPTX файлын оқу функциясы
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_blocks = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_blocks.append(f"--- Слайд {slide_num} ---\n" + "\n".join(slide_text))
    return "\n\n".join(text_blocks)

# PDF файлын оқу функциясы
def extract_text_from_pdf(file):
    text_blocks = []
    with pdfplumber.open(file) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_blocks.append(f"--- Бет {page_num} ---\n{page_text.strip()}")
    return "\n\n".join(text_blocks)

# AI арқылы қатені талдау функциясы
def analyze_content(text, api_key, model, ortho, tech, content):
    client = openai.OpenAI(api_key=api_key)
    
    aspects = []
    if ortho:
        aspects.append("- Орфографиялық, грамматикалық, пунктуациялық қателер мен әріп ауысулары.")
    if tech:
        aspects.append("- Техникалық қателер: Сұрақтардың нөмірленуі, жауап нұсқаларының (A, B, C, D, E) толықтығы, нұсқалардың қайталануы немесе жетіспеуі.")
    if content:
        aspects.append("- Мазмұндық қателер: Сұрақтағы логикалық қайшылықтар, пәндік терминдердің бұрмалануы, дұрыс жауаптың сұраққа сәйкес келмеуі.")

    aspects_str = "\n".join(aspects)

    prompt = f"""
Сен — ҰБТ (Ұлттық бірыңғай тестілеу) материалдары мен тест тапсырмаларын тексеретін жоғары санатты эксперт-методистсің.
Саған тексеруге мынадай мәтін берілді. Мәтінді мұқият талдап, төмендегі критерийлер бойынша қателерді тап:

Критерийлер:
{aspects_str}

Жауапты ОРЫНДАУ ШАРТЫ:
1. Жауапты ТЕК JSON форматында бер. Басқа артық кіріспе немесе қорытынды сөз жазба.
2. Қате табылса, әр қатені мына құрылымда көрсет:
{{
  "errors": [
    {{
      "location": "Қате табылған орын (Слайд/Абзац/Бет нөмірі немесе сұрақ нөмірі)",
      "original": "Қате жазылған фрагмент немесе сөйлем",
      "correction": "Дұрыс нұсқасы / Ұсыныс",
      "category": "Орфографиялық / Техникалық / Мазмұндық",
      "explanation": "Қатенің себебі мен методикалық түсініктемесі"
    }}
  ],
  "summary": "Материалдың жалпы сапасына қысқаша методикалық бағалау"
}}

Егер ешқандай қате табылмаса, "errors" массивін бос қалдыр ([]).

Тексерілетін мәтін:
{text}
"""

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": "Сен JSON форматында ғана жауап беретін сарапшы-методистсің."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
        response_format={"type": "json_object"}
    )
    
    return json.loads(response.choices[0].message.content)

# Файл жүктеу интерфейсі
uploaded_file = st.file_uploader(
    "Файлды таңдаңыз (.docx, .pptx, .pdf):",
    type=["docx", "pptx", "pdf"]
)

if uploaded_file is not None:
    st.info(f"📁 **Жүктелген файл:** {uploaded_file.name} ({round(uploaded_file.size / 1024, 1)} KB)")
    
    extracted_text = ""
    try:
        if uploaded_file.name.endswith(".docx"):
            extracted_text = extract_text_from_docx(uploaded_file)
        elif uploaded_file.name.endswith(".pptx"):
            extracted_text = extract_text_from_pptx(uploaded_file)
        elif uploaded_file.name.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(uploaded_file)
            
        with st.expander("📄 Бөлініп алынған мәтінді көру (Тексеру үшін)"):
            st.text_area("Алынған мәтін:", extracted_text, height=200)
            
    except Exception as e:
        st.error(f"Файлды оқу кезінде қате шықты: {e}")

    # Тексеру батырмасы
    if st.button("🚀 Қателерді тексеруді бастау", type="primary"):
        if not api_key:
            st.warning("⚠️ Өтініш, сол жақтағы баптаулар мәзірінен OpenAI API Key енгізіңіз!")
        elif not extracted_text.strip():
            st.error("❌ Файлдан ешқандай мәтін табылмады.")
        else:
            with st.spinner("🔍 AI методист файлды талдауда... Күте тұрыңыз..."):
                try:
                    result = analyze_content(
                        extracted_text, 
                        api_key, 
                        selected_model, 
                        check_ortho, 
                        check_tech, 
                        check_content
                    )
                    
                    st.subheader("📊 Тексеру Нәтижесі")
                    st.success(f"**Методикалық қорытынды:** {result.get('summary', 'Тексеріс аяқталды.')}")
                    
                    errors = result.get("errors", [])
                    if not errors:
                        st.balloons()
                        st.success("🎉 Тамаша! Материалдан ешқандай қате табылған жоқ.")
                    else:
                        st.error(f"⚠️ Жалпы табылған қателер саны: **{len(errors)}**")
                        
                        # Табылған қателерді блоктарға бөліп шығару
                        for idx, err in enumerate(errors, 1):
                            cat = err.get("category", "Басқа")
                            badge_color = "🔴" if cat == "Мазмұндық" else ("🟡" if cat == "Техникалық" else "🔵")
                            
                            with st.expander(f"{badge_color} #{idx} | [{cat}] {err.get('location', 'Анықталмаған орын')}"):
                                col1, col2 = st.columns(2)
                                with col1:
                                    st.markdown(f"**❌ Файлдағы түпнұсқа:**\n`{err.get('original', '-')}`")
                                with col2:
                                    st.markdown(f"**✅ Дұрыс нұсқасы / Ұсыныс:**\n`{err.get('correction', '-')}`")
                                st.markdown(f"**💡 Түсініктеме:** {err.get('explanation', '-')}")
                                
                except Exception as e:
                    st.error(f"Тексеру кезінде қате орын алды: {e}")